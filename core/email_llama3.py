import os
from django.conf import settings
from groq import Groq
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from better_profanity import profanity
from django.http import JsonResponse, HttpResponse
from docx import Document as DocxDocument
import fitz  # PyMuPDF
import openpyxl
import xlrd

# API Keys
GROQ_SECRET_ACCESS_KEY = settings.GROQ_SECRET_ACCESS_KEY
BHASHINI_API_KEY = settings.BHASHINI_API_KEY
BHASHINI_USER_ID = settings.BHASHINI_USER_ID

# Function to check for inappropriate language
# def contains_inappropriate_language(text: str) -> bool:
#     inappropriate_words = ["stupid", "idiot", "badword3"]
#     return any(word in text.lower() for word in inappropriate_words)

# Function to check for inappropriate language
def contains_inappropriate_language(text: str) -> bool:
    return profanity.contains_profanity(text)

# Function to sanitize input containing inappropriate words
def sanitize_input(input_str):
    return profanity.censor(input_str)

#os.environ["GROQ_SECRET_ACCESS_KEY"] = GROQ_SECRET_ACCESS_KEY


# # Function to generate email
# def generate_email(purpose, num_words, subject, rephrase, to, tone, keywords, contextual_background, call_to_action, additional_details, priority_level, closing_remarks):
#     # Ensure all fields are checked for inappropriate language
#     fields_to_check = [purpose, subject, keywords, contextual_background, call_to_action, additional_details, closing_remarks]
#     if any(contains_inappropriate_language(str(field)) for field in fields_to_check if field is not None):
#         return "Error: Input contains inappropriate language."
    
#     prompt = f"Generate an email of maximum {num_words} words and subject: {subject}, to {to}, maintain a {tone} tone, using the following keywords {keywords}, given the following inputs:"
#     prompt += f"\nPurpose of the mail is {purpose}," if purpose else ""
#     prompt += f"\nConsider the contextual background {contextual_background}," if contextual_background else ""
#     prompt += f"\nWith an expectation of {call_to_action}," if call_to_action else ""
#     prompt += f"\nIncorporate the following additional details: {additional_details}." if additional_details else ""
#     prompt += f"\nThe mail is of {priority_level} priority." if priority_level else ""
#     prompt += f"\nIncorporate the closing remarks {closing_remarks}." if closing_remarks else ""
#     prompt += "\nRephrase the subject" if rephrase == "Y" else ""
    
#     client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)
#     chat_completion = client.chat.completions.create(
#         messages=[{"role": "user", "content": prompt}],
#         model="llama3-70b-8192"
#     )
    
#     return chat_completion.choices[0].message.content


# Function to generate email
# def generate_email(purpose, num_words, subject, rephrase, to, tone, keywords, contextual_background, call_to_action, additional_details, priority_level, closing_remarks):
#     # Ensure all fields are checked for inappropriate language
#     fields_to_check = [purpose, subject, keywords, contextual_background, call_to_action, additional_details, closing_remarks]
#     if any(contains_inappropriate_language(str(field)) for field in fields_to_check if field is not None):
#         return "Error: Input contains inappropriate language."
    
#     prompt = f"Generate an email of maximum {num_words} words and subject: {subject}, to {to}, maintain a {tone} tone, using the following keywords {keywords}, given the following inputs:"
#     prompt += f"\nPurpose of the mail is {purpose}," if purpose else ""
#     prompt += f"\nConsider the contextual background {contextual_background}," if contextual_background else ""
#     prompt += f"\nWith an expectation of {call_to_action}," if call_to_action else ""
#     prompt += f"\nIncorporate the following additional details: {additional_details}." if additional_details else ""
#     prompt += f"\nThe mail is of {priority_level} priority." if priority_level else ""
#     prompt += f"\nIncorporate the closing remarks {closing_remarks}." if closing_remarks else ""
#     prompt += "\nRephrase the subject" if rephrase == "Y" else ""
    
#     client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)
#     chat_completion = client.chat.completions.create(
#         messages=[{"role": "user", "content": prompt}],
#         model="llama3-70b-8192"
#     )
    
#     return chat_completion.choices[0].message.content

def generate_email(purpose, num_words, subject, rephrase, to, tone, keywords, contextual_background, call_to_action, additional_details, priority_level, closing_remarks):
    # Ensure all fields are checked for inappropriate language
    fields_to_check = [purpose, subject, keywords, contextual_background, call_to_action, additional_details, closing_remarks]
    if any(contains_inappropriate_language(str(field)) for field in fields_to_check if field is not None):
        return "Error: Input contains inappropriate language."
    
    prompt = f"Generate an email of maximum {num_words} words and subject: {subject}, to {to}, maintain a {tone} tone, using the following keywords {', '.join(keywords)}, given the following inputs:"
    prompt += f"\nPurpose of the mail is {purpose}," if purpose else ""
    prompt += f"\nConsider the contextual background {contextual_background}," if contextual_background else ""
    prompt += f"\nWith an expectation of {call_to_action}," if call_to_action else ""
    prompt += f"\nIncorporate the following additional details: {additional_details}." if additional_details else ""
    prompt += f"\nThe mail is of {priority_level} priority." if priority_level else ""
    prompt += f"\nIncorporate the closing remarks {closing_remarks}." if closing_remarks else ""
    prompt += "\nRephrase the subject" if rephrase == "Y" else ""
    
    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)
    chat_completion = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.1-70b-versatile",
    )
    
    return chat_completion.choices[0].message.content


def generate_bus_pro(business_intro, proposal_objective, num_words, scope_of_work, project_phases, expected_outcomes, tech_innovations, target_audience, budget_info, timeline, benefits, closing_remarks):
    # Collect all fields to check for inappropriate language
    fields_to_check = [business_intro, proposal_objective, scope_of_work, project_phases, expected_outcomes, tech_innovations, target_audience, budget_info, timeline, benefits, closing_remarks]
    
    # Check if any field contains inappropriate language
    if any(contains_inappropriate_language(str(field)) for field in fields_to_check if field is not None):
        return "Error: Input contains inappropriate language."
    
    # Sanitize input (if needed)
    sanitized_fields = [sanitize_input(str(field)) if field else '' for field in fields_to_check]

    # Reconstruct the prompt with sanitized input
    prompt = f"Generate a business proposal of maximum {num_words} words, given the following inputs: "
    if sanitized_fields[0]:
        prompt += f"Our business details are {sanitized_fields[0]}, "
    if sanitized_fields[1]:
        prompt += f"and the purpose of this proposal is {sanitized_fields[1]}, "
    if sanitized_fields[2]:
        prompt += f"Define the scope of work as {sanitized_fields[2]}. "
    if sanitized_fields[3]:
        prompt += f"The project will be done in the following phases: {sanitized_fields[3]}. "
    if sanitized_fields[4]:
        prompt += f"Reprise the client of these expected outcomes: {sanitized_fields[4]}. "
    if sanitized_fields[5]:
        prompt += f"Mention our following technologies and innovative approaches: {sanitized_fields[5]}. "
    if sanitized_fields[6]:
        prompt += f"Bear in mind that the target audience is: {sanitized_fields[6]}. "
    if sanitized_fields[7]:
        prompt += f"Incorporate this budget info: {sanitized_fields[7]}. "
    if sanitized_fields[8]:
        prompt += f"The timelines we hope to stick to are: {sanitized_fields[8]}. "
    if sanitized_fields[9]:
        prompt += f"Incorporate into the proposal the following benefits: {sanitized_fields[9]}. "
    if sanitized_fields[10]:
        prompt += f"Incorporate the following closing remarks: {sanitized_fields[10]}. "

    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)
    chat_completion = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.1-70b-versatile",
    )

    return chat_completion.choices[0].message.content


def generate_offer_letter(company_details, candidate_name, position_title, department, status, location,
                          start_date, compensation_benefits, work_hours, terms, acceptance_deadline,
                          contact_info, documents_needed, closing_remarks):
    # Collect all fields to check for inappropriate language
    fields_to_check = [
        company_details, candidate_name, position_title, department, status, location,
        start_date, compensation_benefits, work_hours, terms, acceptance_deadline,
        contact_info, documents_needed, closing_remarks
    ]
    
    # Check if any field contains inappropriate language
    inappropriate_key = None
    inappropriate_value = None
    for value in fields_to_check:
        if value and contains_inappropriate_language(value):
            inappropriate_key = value
            inappropriate_value = value
            break

    if inappropriate_key:
        return f"This type of language is not allowed in the input: {inappropriate_value}"

    # Sanitize input fields
    sanitized_fields = [sanitize_input(str(field)) if field else '' for field in fields_to_check]

    # Build the prompt with sanitized inputs
    prompt = "Generate an offer letter given the following inputs: "
    if sanitized_fields[0]:
        prompt += f"\nOur business details are {sanitized_fields[0]}, "
    if sanitized_fields[1]:
        prompt += f"\nCandidate name is {sanitized_fields[1]}, "
    if sanitized_fields[2]:
        prompt += f"\nfor the position of {sanitized_fields[2]}, "
    if sanitized_fields[3]:
        prompt += f"\nin the department: {sanitized_fields[3]}, "
    if sanitized_fields[4]:
        prompt += f"\nas a {sanitized_fields[4]} employee. "
    if sanitized_fields[5]:
        prompt += f"\nExpected to work from: {sanitized_fields[5]}. "
    if sanitized_fields[6]:
        prompt += f"\nCandidate to join on {sanitized_fields[6]}. "
    if sanitized_fields[7]:
        prompt += f"\nCandidate will receive the following compensation and benefits: {sanitized_fields[7]}. "
    if sanitized_fields[8]:
        prompt += f"\nExpected working hours: {sanitized_fields[8]}. "
    if sanitized_fields[9]:
        prompt += f"\nFollowing are the terms of the offer: {sanitized_fields[9]}. "
    if sanitized_fields[10]:
        prompt += f"\nThe last day for accepting the offer is: {sanitized_fields[10]}. "
    if sanitized_fields[11]:
        prompt += f"\nIn case of any queries contact: {sanitized_fields[11]}. "
    if sanitized_fields[12]:
        prompt += f"\nFollowing documents to be produced on the day of joining: {sanitized_fields[12]}. "
    if sanitized_fields[13]:
        prompt += f"\nIncorporate the following closing remarks in the offer letter: {sanitized_fields[13]}. "

    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)

    chat_completion = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": prompt,
            }
        ],
        model="llama-3.1-70b-versatile",
    )

    return chat_completion.choices[0].message.content



def generate_summary(document_context, main_subject, summary_purpose, length_detail, important_elements, audience, tone, format, additional_instructions, document):
    # Collect all fields to check for inappropriate language
    fields_to_check = [
        document_context, main_subject, summary_purpose, length_detail, important_elements, 
        audience, tone, format, additional_instructions, document
    ]
    
    # Check if any field contains inappropriate language
    inappropriate_key = None
    inappropriate_value = None
    for value in fields_to_check:
        if value and contains_inappropriate_language(value):
            inappropriate_key = value
            inappropriate_value = value
            break

    if inappropriate_key:
        return f"This type of language is not allowed in the input: {inappropriate_value}"

    # Sanitize input fields
    sanitized_fields = { 
        "document_context": sanitize_input(str(document_context)) if document_context else '',
        "main_subject": sanitize_input(str(main_subject)) if main_subject else '',
        "summary_purpose": sanitize_input(str(summary_purpose)) if summary_purpose else '',
        "length_detail": sanitize_input(str(length_detail)) if length_detail else '',
        "important_elements": sanitize_input(str(important_elements)) if important_elements else '',
        "audience": sanitize_input(str(audience)) if audience else '',
        "tone": sanitize_input(str(tone)) if tone else '',
        "format": sanitize_input(str(format)) if format else '',
        "additional_instructions": sanitize_input(str(additional_instructions)) if additional_instructions else '',
        "document": sanitize_input(str(document)) if document else ''
    }

    # Build the prompt with sanitized inputs
    prompt = f"Generate a summary of the given document {sanitized_fields['document']} given the following inputs: "
    if sanitized_fields['document_context']:
        prompt += f"\nContext of document: {sanitized_fields['document_context']}, "
    if sanitized_fields['main_subject']:
        prompt += f"\nMain subject: {sanitized_fields['main_subject']}, "
    if sanitized_fields['summary_purpose']:
        prompt += f"\nPurpose of generating summary: {sanitized_fields['summary_purpose']}, "
    if sanitized_fields['length_detail']:
        prompt += f"\nLevel of detail: {sanitized_fields['length_detail']}, "
    if sanitized_fields['important_elements']:
        prompt += f"\nImportant elements: {sanitized_fields['important_elements']}. "
    if sanitized_fields['audience']:
        prompt += f"\nTarget audience: {sanitized_fields['audience']}. "
    if sanitized_fields['tone']:
        prompt += f"\nExpected tone: {sanitized_fields['tone']}. "
    if sanitized_fields['format']:
        prompt += f"\nExpected format: {sanitized_fields['format']}. "
    if sanitized_fields['additional_instructions']:
        prompt += f"\nAdditional Instructions: {sanitized_fields['additional_instructions']}. "

    if not sanitized_fields['document']:
        return "Error: Attach Document!!"

    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)

    chat_completion = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": prompt,
            }
        ],
        model="llama-3.1-70b-versatile",

    )


    return chat_completion.choices[0].message.content




# Function to generate content based on provided parameters
def generate_content(company_info, content_purpose, desired_action, topic_details, keywords, audience_profile, format_structure, num_words, seo_keywords, references):
    inputs = {
        "company_info": company_info,
        "content_purpose": content_purpose,
        "desired_action": desired_action,
        "topic_details": topic_details,
        "keywords": keywords,
        "audience_profile": audience_profile,
        "format_structure": format_structure,
        "seo_keywords": seo_keywords,
        "references": references
    }

    # Check if any input parameter contains inappropriate words
    inappropriate_key = None
    inappropriate_value = None
    for key, value in inputs.items():
        if value and contains_inappropriate_language(value):
            inappropriate_key = key
            inappropriate_value = value
            break

    if inappropriate_key:
        return f"This type of language is not allowed in {inappropriate_key}: '{inappropriate_value}'."

    # Sanitize input fields
    sanitized_inputs = {key: sanitize_input(value) if value else '' for key, value in inputs.items()}

    # Construct the prompt for content generation
    prompt = f"Generate high-quality, engaging content of maximum {num_words} words with the following details:\n"

    if sanitized_inputs['company_info']:
        prompt += f"Company Information: {sanitized_inputs['company_info']}\n"
    if sanitized_inputs['content_purpose']:
        prompt += f"Purpose of Content: {sanitized_inputs['content_purpose']}\n"
    if sanitized_inputs['desired_action']:
        prompt += f"Desired Action: {sanitized_inputs['desired_action']}\n"
    if sanitized_inputs['topic_details']:
        prompt += f"Topic Details: {sanitized_inputs['topic_details']}\n"
    if sanitized_inputs['keywords']:
        prompt += f"Keywords: {sanitized_inputs['keywords']}\n"
    if sanitized_inputs['audience_profile']:
        prompt += f"Audience Profile: {sanitized_inputs['audience_profile']}\n"
    if sanitized_inputs['format_structure']:
        prompt += f"Format and Structure: {sanitized_inputs['format_structure']}\n"
    if sanitized_inputs['seo_keywords']:
        prompt += f"SEO Keywords: {sanitized_inputs['seo_keywords']}\n"
    if sanitized_inputs['references']:
        prompt += f"References to Cite: {sanitized_inputs['references']}\n"

    # Additional instructions for the content creation
    prompt += (
        "\nInstructions:\n"
        "- Ensure the content is engaging, informative, and relevant to the specified audience.\n"
        "- Highlight the benefits and unique aspects of the topic to capture the audience's interest.\n"
        "- Use a professional tone and clear language to communicate effectively.\n"
        "- Incorporate the provided keywords naturally and strategically for SEO optimization.\n"
        "- Maintain accuracy and avoid any hallucinations or false information.\n"
        "- Adhere to the specified format and structure to meet the content requirements.\n"
    )

    # Generate content using Groq API
    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)

    chat_completion = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": prompt,
            }
        ],
        model="llama-3.1-70b-versatile",
    )

    return chat_completion.choices[0].message.content


def generate_sales_script(company_details, num_words, product_descriptions, features_benefits, pricing_info, promotions, target_audience, sales_objectives,
                          competitive_advantage, compliance):
    inputs = {
        "Company Details": company_details,
        "Product Descriptions": product_descriptions,
        "Features and Benefits": features_benefits,
        "Pricing Info": pricing_info,
        "Promotions": promotions,
        "Target Audience": target_audience,
        "Sales Objectives": sales_objectives,
        "Competitive Advantage": competitive_advantage,
        "Compliance": compliance,
        "Number Of Words": num_words
    }

    # Check if any input parameter contains inappropriate words
    inappropriate_key = None
    inappropriate_value = None
    for key, value in inputs.items():
        if value and contains_inappropriate_language(value):
            inappropriate_key = key
            inappropriate_value = value
            break

    if inappropriate_key:
        return f"This type of language is not allowed in {inappropriate_key}: {inappropriate_value}"

    # Sanitize input fields
    sanitized_inputs = {key: sanitize_input(value) if value else '' for key, value in inputs.items()}

    # Build the prompt for generating the sales script
    prompt = f"Generate a sales script of maximum {sanitized_inputs['Number Of Words']} words, given the following inputs: "
    for key, value in sanitized_inputs.items():
        if key != 'Number Of Words' and value:
            prompt += f"\n- {key}: {value}"

    prompt += (
        "\n\nInstructions:\n"
        "- Ensure the script is professional and persuasive.\n"
        "- Avoid any hallucinations or fabricated information. Use only the provided details.\n"
        "- Maintain accuracy and factual integrity throughout the script.\n"
        "- Avoid using any inappropriate words or foul language."
    )

    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)

    chat_completion = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": prompt,
            }
        ],
        model="llama-3.1-70b-versatile",
    )

    return chat_completion.choices[0].message.content


def bhashini_translate(text: str,  to_code: str = "Hindi", from_code: str = "English",user_id: str=BHASHINI_USER_ID, api_key: str=BHASHINI_API_KEY) -> dict:
    """Translates text from source language to target language using the Bhashini API.

    Args:
        text (str): The text to translate.
        from_code (str): Source language code. Default is 'en' (English).
        to_code (str): Target language code. Default is 'te' (Telugu).
        user_id (str): User ID for the API.
        api_key (str): API key for authentication.

    Returns:
        dict: A dictionary with the status code, message, and translated text or error info.
    """
    lang_dict = {
        "English" :"en",
        "Hindi" :"hi",
        "Tamil" :"ta",
        "Telugu": "te",
        "Marathi": "mr",
        "Kannada" : "kn",
        "Bengali" : "bn",
        "Odia": "or",
        "Assamese": "as",
        "Punjabi": "pa",
        "Malayalam": "ml",
        "Gujarati": "gu",
        "Urdu" : "ur",
        "Sanskrit" : "sa",
        "Nepali" : "ne",
        "Bodo": "brx",
        "Maithili" : "mai",
        "Sindhi" : "sd",
        "Tamil" : "ta",

    }

    print(33333, text, 4444, to_code)
    from_code = lang_dict[from_code]
    to_code = lang_dict[to_code]
    print(222222, from_code, to_code)



    # Setup the initial request to get model configurations
    url = 'https://meity-auth.ulcacontrib.org/ulca/apis/v0/model/getModelsPipeline'
    headers = {
        "Content-Type": "application/json",
        "userID": user_id,
        "ulcaApiKey": api_key
    }
    payload = {
        "pipelineTasks": [{"taskType": "translation", "config": {"language": {"sourceLanguage": from_code, "targetLanguage": to_code}}}],
        "pipelineRequestConfig": {"pipelineId": "64392f96daac500b55c543cd"}
    }
    response = requests.post(url, json=payload, headers=headers)

    if response.status_code != 200:
        return {"status_code": response.status_code, "message": "Error in translation request", "translated_content": None}

    # Process the response to setup the translation execution
    response_data = response.json()
    service_id = response_data["pipelineResponseConfig"][0]["config"][0]["serviceId"]
    callback_url = response_data["pipelineInferenceAPIEndPoint"]["callbackUrl"]
    headers2 = {
        "Content-Type": "application/json",
        response_data["pipelineInferenceAPIEndPoint"]["inferenceApiKey"]["name"]: response_data["pipelineInferenceAPIEndPoint"]["inferenceApiKey"]["value"]
    }
    compute_payload = {
        "pipelineTasks": [{"taskType": "translation", "config": {"language": {"sourceLanguage": from_code, "targetLanguage": to_code}, "serviceId": service_id}}],
        "inputData": {"input": [{"source": text}], "audio": [{"audioContent": None}]}
    }

    # Execute the translation
    compute_response = requests.post(callback_url, json=compute_payload, headers=headers2)
    if compute_response.status_code != 200:
        return {"status_code": compute_response.status_code, "message": "Error in translation", "translated_content": None}

    compute_response_data = compute_response.json()
    translated_content = compute_response_data["pipelineResponse"][0]["output"][0]["target"]

    return {"status_code": 200, "message": "Translation successful", "translated_content": translated_content}



def generate_slide_content(st, title, special_instructions, document_content=None):
    prompt = f"Generate content to be put in ppt slide with title {st}. The overall subject of the presentation is {title}. "
    if special_instructions:
        prompt += f"\nPay attention to the following points: {special_instructions} while generating content. "
    if document_content:
        prompt += f"\nUse the following document content as a reference: {document_content[:2000]}... "
    prompt += (
        "Content should be in the form of points and should be just sufficient to fit on a single slide {highest priority} "
        "Output only the slide contents - avoid any text describing the content. "
        "Do not include title in the content. No bold text. Avoid the text 'Here is the content for the PPT slide'. "
        "Do not include the bullet point symbols."
    )

    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)
    chat_completion = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.1-70b-versatile",

    )
    slide_content = chat_completion.choices[0].message.content
    return slide_content

def generate_slide_titles(title, num_slides, special_instructions):
    prompt = f"Generate titles for {num_slides} slides on the subject {title}. "
    if special_instructions:
        prompt += f"\nPay attention to the following points: {special_instructions} while generating titles. "
    prompt += ("Titles should be in the form of a python list object with {num_slides} elements. Output only the list object without any text before or after the list."
            "The Title should not exceed more than two lines."
            "Do not Hallucinate")

    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)
    chat_completion = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.1-70b-versatile",
    )
    title_list = chat_completion.choices[0].message.content
    return title_list

def add_slide(prs, title, content, bg_image):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_text_frame = title_placeholder.text_frame
    title_font_size = Pt(32)  # Title font size
    small_font_size = Pt(20)
    title_text_frame.clear()  # Clear any existing paragraphs

    p = title_text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title.split("(contd.)")[0]
    run.font.size = title_font_size
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    if "(contd.)" in title:
        run = p.add_run()
        run.text = "(contd.)"
        run.font.size = small_font_size
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 51, 102)

    p.alignment = PP_ALIGN.CENTER

    content_text_frame = content_placeholder.text_frame
    content_font_size = Pt(20)  # Updated default content font size
    content_text_frame.clear()  # Clear any existing paragraphs

    for point in content:
        point = point.lstrip("*•")  # Remove leading bullet symbols
        p = content_text_frame.add_paragraph()
        p.text = point.strip()
        p.font.size = content_font_size
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.LEFT

    while not check_text_fit(prs, content_text_frame):
        content_font_size -= Pt(2)
        for paragraph in content_text_frame.paragraphs:
            if paragraph.font.size:
                paragraph.font.size = content_font_size

    # Set background image
    if bg_image is None:
        bg_image = settings.DEFAULT_BACKGROUND_IMAGE_PATH  # Path to the default background image

    left = top = Inches(0)
    pic = slide.shapes.add_picture(bg_image, left, top, width=prs.slide_width, height=prs.slide_height)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

def check_text_fit(prs, text_frame):
    slide_height = prs.slide_height
    total_height = sum((paragraph.font.size.pt if paragraph.font.size else Pt(18).pt) * len(paragraph.text.split('\n')) for paragraph in text_frame.paragraphs)
    return total_height <= slide_height

def create_presentation(title, num_slides, special_instructions, bg_image):
    prs = Presentation()
    slide_titles = generate_slide_titles(title, num_slides, special_instructions)
    slide_titles = slide_titles.replace('[', '').replace(']', '').replace('"', '').split(',')

    max_points_per_slide = 4     # Adjust this value based on how much content you want per slide

    total_slides_generated = 0

    for st in slide_titles:
        if total_slides_generated >= num_slides:
            break

        slide_content = generate_slide_content(st, title, special_instructions).replace("*", '').split('\n')
        current_content = []
        slide_count = 1

        for point in slide_content:
            current_content.append(point.strip())
            if len(current_content) >= max_points_per_slide:
                add_slide(prs, st if slide_count == 1 else f"{st} (contd.)", current_content, bg_image)
                current_content = []
                slide_count += 1
                total_slides_generated += 1

                if total_slides_generated >= num_slides:
                    break

        if current_content and total_slides_generated < num_slides:
            add_slide(prs, st if slide_count == 1 else f"{st} (contd.)", current_content, bg_image)
            total_slides_generated += 1

    return prs







def extract_document_content(file_path):
    if file_path.endswith('.docx'):
        doc = DocxDocument(file_path)
        return '\n'.join([para.text for para in doc.paragraphs])
    elif file_path.endswith('.pdf'):
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        return text
    elif file_path.endswith('.xlsx'):
        wb = openpyxl.load_workbook(file_path)
        sheet = wb.active
        text = ""
        for row in sheet.iter_rows(values_only=True):
            text += ' '.join([str(cell) for cell in row if cell is not None]) + '\n'
        return text
    elif file_path.endswith('.xls'):
        wb = xlrd.open_workbook(file_path)
        sheet = wb.sheet_by_index(0)
        text = ""
        for row_idx in range(sheet.nrows):
            row = sheet.row(row_idx)
            text += ' '.join([str(cell.value) for cell in row if cell.value]) + '\n'
        return text
    else:
        raise ValueError("Unsupported file type")








