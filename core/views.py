from io import BytesIO
import os
import random
from django.shortcuts import render
from django.contrib.auth import authenticate, login, logout
from django.contrib.auth.decorators import login_required
from django.http import HttpResponse, JsonResponse
from .email_llama3 import add_slide, create_presentation, extract_document_content, generate_email, bhashini_translate,generate_bus_pro, generate_offer_letter, generate_slide_content, generate_slide_titles, generate_summary, generate_content, generate_sales_script  
from django.contrib.auth.models import User
from django.utils.crypto import get_random_string
from django.utils import timezone
from django.conf import settings
from datetime import timedelta
from .models import PasswordResetRequest, Profile
from django.core.mail import send_mail, BadHeaderError
from django.contrib.auth import update_session_auth_hash
from django.views.decorators.csrf import csrf_exempt
import json
from rest_framework_simplejwt.tokens import RefreshToken
from django.contrib.auth import authenticate
from rest_framework_simplejwt.tokens import RefreshToken
from rest_framework.decorators import api_view, permission_classes
from rest_framework.permissions import IsAuthenticated
from django.contrib.auth.password_validation import validate_password  
import json
from Crypto.Cipher import AES
from Crypto.Util.Padding import pad, unpad
import base64
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.contrib.auth.models import User
from django.core.exceptions import ValidationError
from rest_framework import status
from rest_framework.decorators import api_view, permission_classes
from rest_framework.renderers import BaseRenderer
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from django.core.files.storage import default_storage
from django.utils.dateparse import parse_date
from django.core.files.uploadedfile import InMemoryUploadedFile
from django.core.files.storage import default_storage
from django.shortcuts import render
import fitz  # PyMuPDF
from docx import Document as DocxDocument
import logging

logger = logging.getLogger(__name__)


def test_report(request):
    report_path = os.path.join(os.path.dirname(__file__), 'report.html')
    with open(report_path, 'r') as file:
        report_content = file.read()
    return render(request, 'test_report.html', {'report_content': report_content})


# Base64-encoded AES IV and Secret Key
AES_IV_b64 = settings.AES_IV
AES_SECRET_KEY_b64 = settings.AES_SECRET_KEY
ENCRYPTION_IV_b64 = settings.ENCRYPTION_IV
ENCRYPTION_SECRET_KEY_b64 = settings.ENCRYPTION_SECRET_KEY

# Decode Base64 strings to bytes
AES_IV = base64.b64decode(AES_IV_b64)
AES_SECRET_KEY = base64.b64decode(AES_SECRET_KEY_b64)
ENCRYPTION_IV = base64.b64decode(ENCRYPTION_IV_b64)
ENCRYPTION_SECRET_KEY = base64.b64decode(ENCRYPTION_SECRET_KEY_b64)

# Decode Base64 strings to bytes
AES_IV = base64.b64decode(AES_IV_b64)
AES_SECRET_KEY = base64.b64decode(AES_SECRET_KEY_b64)

# Ensure IV is 16 bytes long (128 bits)
if len(AES_IV) != 16:
    raise ValueError("AES IV must be 16 bytes long")


class CustomAesRenderer(BaseRenderer):
    media_type = 'application/octet-stream'
    format = 'aes'

    def render(self, data, media_type=None, renderer_context=None):
        plaintext = json.dumps(data)
        padded_plaintext = pad(plaintext.encode(), 16)
        cipher = AES.new(AES_SECRET_KEY, AES.MODE_CBC, AES_IV)
        ciphertext = cipher.encrypt(padded_plaintext)
        ciphertext_b64 = base64.b64encode(ciphertext).decode()
        response = {'ciphertext': ciphertext_b64}
        return json.dumps(response)

def landing(request):
    return render(request, 'landing.html')

def about(request):
    return render(request, 'about.html')

def encrypt_data(data):
    plaintext = json.dumps(data)
    padded_plaintext = pad(plaintext.encode(), 16)
    cipher = AES.new(ENCRYPTION_SECRET_KEY, AES.MODE_CBC, ENCRYPTION_IV)
    ciphertext = cipher.encrypt(padded_plaintext)
    ciphertext_b64 = base64.b64encode(ciphertext).decode()
    return ciphertext_b64

def decrypt_data(encrypted_data):
    try:
        encrypted_data_bytes = base64.b64decode(encrypted_data)
        cipher = AES.new(ENCRYPTION_SECRET_KEY, AES.MODE_CBC, ENCRYPTION_IV)
        decrypted_data = unpad(cipher.decrypt(encrypted_data_bytes), 16)
        return decrypted_data.decode('utf-8')
    except Exception as e:
        raise ValueError(f"Decryption error: {e}")

# @csrf_exempt
# def add_user(request):
#     try:
#         logger.debug("Request received")
#         # Extract and decrypt the incoming payload
#         encrypted_content = json.loads(request.body.decode('utf-8')).get('encrypted_content')
#         if not encrypted_content:
#             logger.error("No encrypted content found in the request.")
#             return JsonResponse({'error': 'No encrypted content found in the request.'}, status=400)
        
#         decrypted_content = decrypt_data(encrypted_content)
#         data = json.loads(decrypted_content)
#         logger.debug("Request body: %s", data)

#         username = data.get('username')
#         email = data.get('email')
#         password = data.get('password')
#         confirm_password = data.get('confirm_password')

#         if not username or not email or not password or not confirm_password:
#             logger.error("Missing field(s)")
#             return JsonResponse({'error': 'All fields are required.'}, status=400)

#         if password != confirm_password:
#             logger.error("Passwords do not match")
#             return JsonResponse({'error': 'Passwords do not match.'}, status=400)

#         try:
#             validate_password(password)
#             logger.debug("Password validation passed")
#         except ValidationError as e:
#             logger.error("Password validation error: %s", e.messages)
#             return JsonResponse({'error': list(e.messages)}, status=400)

#         if User.objects.filter(username=username).exists():
#             logger.error("Username is already taken")
#             return JsonResponse({'error': 'Username is already taken.'}, status=400)

#         if User.objects.filter(email=email).exists():
#             logger.error("Email is already registered")
#             return JsonResponse({'error': 'Email is already registered.'}, status=400)

#         user = User.objects.create_user(username=username, email=email, password=password)
#         logger.debug("User created: %s", user)

#         # Encrypt the response content
#         encrypted_response = encrypt_data({'success': 'User created successfully.'})

#         return JsonResponse({'encrypted_content': encrypted_response}, status=201)

#     except json.JSONDecodeError:
#         logger.error("Invalid JSON format")
#         return JsonResponse({'error': 'Invalid JSON format.'}, status=400)
#     except Exception as e:
#         logger.exception("Unexpected error")
#         return JsonResponse({'error': str(e)}, status=500)
    
@csrf_exempt
def add_user(request):
    try:
        logger.debug("Request received")
        # Extract and decrypt the incoming payload
        encrypted_content = json.loads(request.body.decode('utf-8')).get('encrypted_content')
        if not encrypted_content:
            logger.error("No encrypted content found in the request.")
            return JsonResponse({'error': 'No encrypted content found in the request.'}, status=400)
        
        decrypted_content = decrypt_data(encrypted_content)
        data = json.loads(decrypted_content)
        logger.debug("Request body: %s", data)

        first_name = data.get('first_name')
        last_name = data.get('last_name')
        username = data.get('username')
        email = data.get('email')
        password = data.get('password')
        confirm_password = data.get('confirm_password')

        if not first_name or not last_name or not username or not email or not password or not confirm_password:
            logger.error("Missing field(s)")
            return JsonResponse({'error': 'All fields are required.'}, status=400)

        if password != confirm_password:
            logger.error("Passwords do not match")
            return JsonResponse({'error': 'Passwords do not match.'}, status=400)

        try:
            validate_password(password)
            logger.debug("Password validation passed")
        except ValidationError as e:
            logger.error("Password validation error: %s", e.messages)
            return JsonResponse({'error': list(e.messages)}, status=400)

        if User.objects.filter(username=username).exists():
            logger.error("Username is already taken")
            return JsonResponse({'error': 'Username is already taken.'}, status=400)

        if User.objects.filter(email=email).exists():
            logger.error("Email is already registered")
            return JsonResponse({'error': 'Email is already registered.'}, status=400)

        user = User.objects.create_user(username=username, email=email, password=password, first_name=first_name, last_name=last_name)
        logger.debug("User created: %s", user)

        # Encrypt the response content
        encrypted_response = encrypt_data({'success': 'User created successfully.'})

        return JsonResponse({'encrypted_content': encrypted_response}, status=201)

    except json.JSONDecodeError:
        logger.error("Invalid JSON format")
        return JsonResponse({'error': 'Invalid JSON format.'}, status=400)
    except Exception as e:
        logger.exception("Unexpected error")
        return JsonResponse({'error': str(e)}, status=500)



@csrf_exempt
def generate_otp():
    return ''.join(random.choices('0123456789', k=6))

@csrf_exempt
def send_otp(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            email = data.get('email')
            logger.debug(f"Received OTP request for email: {email}")

            try:
                user = User.objects.get(email=email)
            except User.DoesNotExist:
                logger.warning(f"Email does not exist: {email}")
                return JsonResponse({'error': 'Email does not exist'}, status=404)

            # Generate OTP
            otp = generate_otp()
            expiry_time = timezone.now() + timedelta(minutes=10)

            # Store OTP and expiry time in the database
            PasswordResetRequest.objects.update_or_create(
                user=user,
                defaults={
                    'otp': otp,
                    'expiry_time': expiry_time
                }
            )
            logger.info(f"Generated OTP for user {user.username}")

            # Send OTP via email
            subject = 'Password Reset OTP'
            message = f'Your OTP for password reset is {otp}. This OTP is valid only for 10 minutes.'
            from_email = settings.DEFAULT_FROM_EMAIL
            recipient_list = [email]

            try:
                send_mail(subject, message, from_email, recipient_list, fail_silently=False)
                logger.info(f"OTP email sent to {email}")
            except Exception as e:
                logger.error(f"Error sending email: {str(e)}")
                return JsonResponse({'error': f'Error sending email: {str(e)}'}, status=500)

            return JsonResponse({'success': 'OTP sent successfully'}, status=200)

        except json.JSONDecodeError:
            logger.error("Invalid JSON format in request")
            return JsonResponse({'error': 'Invalid JSON format'}, status=400)
        except Exception as e:
            logger.error(f"Unexpected error: {str(e)}")
            return JsonResponse({'error': str(e)}, status=500)

    logger.error("Invalid request method")
    return JsonResponse({'error': 'Invalid request method'}, status=405)


@csrf_exempt
@api_view(['POST'])
def forgot_password(request):
    try:
        data = json.loads(request.body)
        email = data.get('email')
        logger.debug(f"Received password reset request for email: {email}")

        try:
            user = User.objects.get(email=email)
        except User.DoesNotExist:
            logger.warning(f"Email does not exist: {email}")
            return JsonResponse({'error': 'Email does not exist'}, status=404)

        # Generate OTP
        otp = generate_otp()
        expiry_time = timezone.now() + timedelta(minutes=10)
        logger.info(f"Generated OTP for user {user.username}")

        # Store OTP and expiry time in the database
        PasswordResetRequest.objects.create(user=user, otp=otp, expiry_time=expiry_time)
        logger.debug(f"Stored OTP and expiry time in the database for user {user.username}")

        # Send OTP via email
        subject = 'Password Reset OTP'
        message = f'Your OTP for password reset is {otp}. This OTP is valid only for 10 minutes.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [email]

        try:
            send_mail(subject, message, from_email, recipient_list, fail_silently=False)
            logger.info(f"OTP email sent to {email}")
        except BadHeaderError:
            logger.error("Invalid header found.")
            return JsonResponse({'error': 'Invalid header found.'}, status=400)
        except Exception as e:
            logger.error(f"Error sending email: {str(e)}")
            return JsonResponse({'error': f'Error sending email: {str(e)}'}, status=500)

        return JsonResponse({'success': 'OTP sent successfully'}, status=200)

    except json.JSONDecodeError:
        logger.error("Invalid JSON format in request")
        return JsonResponse({'error': 'Invalid JSON format'}, status=400)
    except Exception as e:
        logger.error(f"Unexpected error: {str(e)}")
        return JsonResponse({'error': str(e)}, status=500)


@csrf_exempt
def signin(request):
    if request.method == 'POST':
        try:
            # Extract and decrypt the incoming payload
            encrypted_content = json.loads(request.body.decode('utf-8')).get('encrypted_content')
            if not encrypted_content:
                logger.warning('No encrypted content found in the request.')
                return JsonResponse({'error': 'No encrypted content found in the request.'}, status=400)
            
            decrypted_content = decrypt_data(encrypted_content)
            data = json.loads(decrypted_content)
            logger.debug(f'Decrypted content: {data}')

            email = data.get('email')
            password = data.get('password')
            
            if not email or not password:
                logger.warning('Email and password are required')
                return JsonResponse({'error': 'Email and password are required'}, status=400)

            try:
                user = User.objects.get(email=email)
            except User.DoesNotExist:
                logger.warning('Invalid email or password')
                return JsonResponse({'error': 'Invalid email or password'}, status=400)

            user = authenticate(request, username=user.username, password=password)
            if user is not None:
                login(request, user)
                refresh = RefreshToken.for_user(user)
                logger.info(f'User {user.username} authenticated successfully')

                # Encrypt the response content
                encrypted_response = encrypt_data({
                    'success': 'User authenticated',
                    'access': str(refresh.access_token),
                    'refresh': str(refresh)
                })

                return JsonResponse({'encrypted_content': encrypted_response}, status=200)
            else:
                logger.warning('Invalid email or password')
                return JsonResponse({'error': 'Invalid email or password'}, status=400)
        except json.JSONDecodeError:
            logger.error('Invalid JSON format in request')
            return JsonResponse({'error': 'Invalid JSON'}, status=400)
        except Exception as e:
            logger.error(f'Internal server error: {str(e)}')
            return JsonResponse({'error': 'Internal server error'}, status=500)
    else:
        logger.error('Invalid request method')
        return JsonResponse({'error': 'Invalid request method'}, status=405)


@csrf_exempt
def reset_password(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            logger.debug(f'Request data: {data}')
        except json.JSONDecodeError:
            logger.error('Invalid JSON format in request')
            return JsonResponse({'error': 'Invalid JSON'}, status=400)

        email = data.get('email')
        otp = data.get('otp')
        new_password = data.get('new_password')
        confirm_password = data.get('confirm_password')

        if not all([email, otp, new_password, confirm_password]):
            logger.warning('All fields are required')
            return JsonResponse({'error': 'All fields are required'}, status=400)

        if new_password != confirm_password:
            logger.warning('Passwords do not match')
            return JsonResponse({'error': 'Passwords do not match'}, status=400)

        try:
            user = User.objects.get(email=email)
            logger.info(f'User found: {user.username}')
        except User.DoesNotExist:
            logger.warning(f'User with email {email} does not exist')
            return JsonResponse({'error': 'User with this email does not exist'}, status=404)

        try:
            password_reset_request = PasswordResetRequest.objects.get(user=user, otp=otp)
            logger.info('Password reset request found')
        except PasswordResetRequest.DoesNotExist:
            logger.warning('Invalid OTP')
            return JsonResponse({'error': 'Invalid OTP'}, status=400)

        if password_reset_request.expiry_time < timezone.now():
            logger.warning('OTP has expired')
            return JsonResponse({'error': 'OTP has expired'}, status=400)

        user.set_password(new_password)
        user.save()
        logger.info(f'Password for user {user.username} reset successfully')

        password_reset_request.delete()
        logger.info('Password reset request deleted')

        return JsonResponse({'success': 'Password reset successfully'}, status=200)
    else:
        logger.error('Invalid request method')
        return JsonResponse({'error': 'Invalid request method'}, status=405)

@csrf_exempt
@api_view(['POST'])
@permission_classes([IsAuthenticated])
def email_generator(request):
    if request.method == 'POST':
        try:
            # Extract and decrypt the incoming payload
            encrypted_content = json.loads(request.body.decode('utf-8')).get('encrypted_content')
            logger.debug(f'Encrypted content received: {encrypted_content}')

            if not encrypted_content:
                logger.warning('No encrypted content found in the request.')
                return JsonResponse({'error': 'No encrypted content found in the request.'}, status=400)

            decrypted_content = decrypt_data(encrypted_content)
            logger.debug(f'Decrypted content: {decrypted_content}')
            data = json.loads(decrypted_content)

            purpose = data.get('purpose')
            if purpose == 'others':
                purpose = data.get('otherPurpose')
            num_words = data.get('num_words')
            subject = data.get('subject')
            rephrase = data.get('rephraseSubject', False)
            to = data.get('to')
            tone = data.get('tone')
            keywords_str = data.get('keywords', '')  # Fetch the keywords as a comma-separated string
            keywords = [keyword.strip() for keyword in keywords_str.split(',')] if keywords_str else []
            contextual_background = data.get('contextualBackground')
            call_to_action = data.get('callToAction')
            if call_to_action == 'Other':
                call_to_action = data.get('otherCallToAction')
            additional_details = data.get('additionalDetails')
            priority_level = data.get('priorityLevel')
            closing_remarks = data.get('closingRemarks')

            logger.info(f'Generating email with the following data: {data}')

            generated_content = generate_email(
                purpose, num_words, subject, rephrase, to, tone, keywords,
                contextual_background, call_to_action, additional_details,
                priority_level, closing_remarks
            )

            if generated_content:
                logger.info('Email content generated successfully.')
                # Encrypt the response content
                encrypted_response = encrypt_data({'generated_content': generated_content})
                logger.debug(f'Encrypted response: {encrypted_response}')

                return JsonResponse({'encrypted_content': encrypted_response}, status=200)

            logger.error('Failed to generate email.')
            return JsonResponse({'error': 'Failed to generate email. Please try again.'}, status=500)

        except json.JSONDecodeError:
            logger.error('Invalid JSON format in request.')
            return JsonResponse({'error': 'Invalid JSON format. Please provide valid JSON data.'}, status=400)
        except ValueError as e:
            logger.warning(f'ValueError: {str(e)}')
            return JsonResponse({'error': str(e)}, status=400)
        except Exception as e:
            logger.error(f'Unexpected error: {str(e)}')
            return JsonResponse({'error': str(e)}, status=500)

    logger.error('Method not allowed.')
    return JsonResponse({'error': 'Method not allowed.'}, status=405)

@csrf_exempt
@api_view(['POST'])
@permission_classes([IsAuthenticated])
def translate_content(request):
    translated_content = None
    error = None
    language = ""

    if request.method == 'POST':
        try:
            data = request.data  # Assuming request body is in JSON format
            logger.debug(f'Request data received: {data}')

            generated_content = data.get('generated_content')
            language = data.get('language')

            if not generated_content or not language:
                logger.warning('Both generated_content and language are required fields.')
                return JsonResponse({'error': 'Both generated_content and language are required fields.'}, status=400)

            logger.info(f'Translating content: {generated_content} to language: {language}')
            response = bhashini_translate(generated_content, language)

            if response["status_code"] == 200:
                translated_content = response["translated_content"]
                logger.info('Content translated successfully.')
                return JsonResponse({
                    'generated_content': generated_content,
                    'translated_content': translated_content,
                    'selected_language': language
                }, status=200)
            else:
                logger.error('Translation failed with status code: {}'.format(response["status_code"]))
                return JsonResponse({'error': 'Translation failed.'}, status=500)

        except Exception as e:
            logger.error(f'Unexpected error: {str(e)}')
            return JsonResponse({'error': str(e)}, status=500)

    logger.error('Method not allowed.')
    return JsonResponse({'error': 'Method not allowed.'}, status=405)


def translate(request):
    translated_text = None
    error = None
    input_text = ""
    from_language = ""
    to_language = ""

    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            logger.debug(f'Request data received: {data}')
            
            input_text = data.get('input_text', '')
            from_language = data.get('from_language', '')
            to_language = data.get('to_language', '')

            if input_text and from_language and to_language:
                try:
                    logger.info(f'Translating text from {from_language} to {to_language}')
                    
                    # Replace this with your translation function call
                    translated_text = bhashini_translate(input_text, to_language, from_language)
                    translated_text = translated_text["translated_content"]
                    
                    logger.info('Translation successful')
                    logger.debug(f'Input text: {input_text}')
                    logger.debug(f'Translated text: {translated_text}')
                except Exception as e:
                    error = f"Error during translation: {str(e)}"
                    logger.error(error)
            else:
                error = "Please provide the input text and select both languages."
                logger.warning(error)
        except json.JSONDecodeError:
            error = "Invalid JSON format received."
            logger.error(error)
    else:
        error = 'Invalid request method'
        logger.warning(error)

    response = {
        'translated_text': translated_text,
        'error': error,
        'input_text': input_text,
        'from_language': from_language,
        'to_language': to_language
    }
    
    logger.debug(f'Response: {response}')
    return JsonResponse(response)


@csrf_exempt
@api_view(['POST'])
@permission_classes([IsAuthenticated])
def business_proposal_generator(request):
    if request.method == 'POST':
        try:
            # Extract and decrypt the incoming payload
            encrypted_content = json.loads(request.body.decode('utf-8')).get('encrypted_content')
            if not encrypted_content:
                logger.warning('No encrypted content found in the request.')
                return JsonResponse({'error': 'No encrypted content found in the request.'}, status=400)

            decrypted_content = decrypt_data(encrypted_content)
            data = json.loads(decrypted_content)
            logger.debug(f'Decrypted content: {data}')

            business_intro = data.get('businessIntroduction')
            proposal_objective = data.get('proposalObjective')
            num_words = data.get('numberOfWords')
            scope_of_work = data.get('scopeOfWork')
            project_phases = data.get('projectPhases')
            expected_outcomes = data.get('expectedOutcomes')
            tech_innovations = data.get('technologiesAndInnovations')  # Combined field
            target_audience = data.get('targetAudience')
            budget_info = data.get('budgetInformation')
            timeline = data.get('timeline')
            benefits = data.get('benefitsToRecipient')
            closing_remarks = data.get('closingRemarks')

            logger.info('Generating business proposal content.')
            proposal_content = generate_bus_pro(
                business_intro, proposal_objective, num_words, scope_of_work,
                project_phases, expected_outcomes, tech_innovations, target_audience,
                budget_info, timeline, benefits, closing_remarks
            )

            # Encrypt the response content
            encrypted_content = encrypt_data({'generated_content': proposal_content})
            logger.info('Business proposal content generated successfully.')

            return JsonResponse({'encrypted_content': encrypted_content}, status=200)

        except json.JSONDecodeError:
            logger.error('Invalid JSON format received.')
            return JsonResponse({'error': 'Invalid JSON format. Please provide valid JSON data.'}, status=400)
        except ValueError as e:
            logger.error(f'ValueError: {str(e)}')
            return JsonResponse({'error': str(e)}, status=400)
        except Exception as e:
            logger.error(f'Exception: {str(e)}')
            return JsonResponse({'error': str(e)}, status=500)

    logger.warning('Method not allowed.')
    return JsonResponse({'error': 'Method not allowed.'}, status=405)


@csrf_exempt
@api_view(['GET'])
@permission_classes([IsAuthenticated])
def fetch_first_name(request):
    user = request.user
    logger.debug(f"User {user.username} accessed the fetch_first_name view.")

    response_data = {
        'first_name': user.first_name,
    }

    logger.debug(f"Returning response data: {response_data}")
    return JsonResponse(response_data)

@csrf_exempt
@api_view(['POST'])
@permission_classes([IsAuthenticated])
def offer_letter_generator(request):
    try:
        encrypted_content = json.loads(request.body.decode('utf-8')).get('encrypted_content')
        if not encrypted_content:
            logger.warning('No encrypted content found in the request.')
            return JsonResponse({'error': 'No encrypted content found in the request.'}, status=400)

        decrypted_content = decrypt_data(encrypted_content)
        data = json.loads(decrypted_content)
        logger.debug(f'Decrypted content: {data}')

        company_details = data.get('companyDetails')
        candidate_name = data.get('candidateFullName')
        position_title = data.get('positionTitle')
        department = data.get('department')
        status = data.get('status')
        location = data.get('location')
        start_date = data.get('expectedStartDate')
        compensation_benefits = data.get('compensationBenefits')  # Merged field
        work_hours = data.get('workHours')
        terms = data.get('termsConditions')
        acceptance_deadline = data.get('deadline')
        contact_info = data.get('contactInfo')
        documents_needed = data.get('documentsNeeded')
        closing_remarks = data.get('closingRemarks')

        logger.info('Generating offer letter content.')
        offer_letter_content = generate_offer_letter(
            company_details,  candidate_name, position_title, department, status,
            location, start_date, compensation_benefits, work_hours,
            terms, acceptance_deadline, contact_info, documents_needed, closing_remarks
        )

        if offer_letter_content:
            encrypted_content = encrypt_data({'generated_content': offer_letter_content})
            logger.info('Offer letter content generated successfully.')
            return JsonResponse({'encrypted_content': encrypted_content}, status=200)

        logger.error('Failed to generate offer letter content.')
        return JsonResponse({'error': 'Failed to generate offer letter. Please try again.'}, status=500)

    except json.JSONDecodeError:
        logger.error('Invalid JSON format received.')
        return JsonResponse({'error': 'Invalid JSON format. Please provide valid JSON data.'}, status=400)
    except ValueError as e:
        logger.error(f'ValueError: {str(e)}')
        return JsonResponse({'error': str(e)}, status=400)
    except Exception as e:
        logger.error(f'Exception: {str(e)}')
        return JsonResponse({'error': str(e)}, status=500)

    logger.warning('Method not allowed.')
    return JsonResponse({'error': 'Method not allowed.'}, status=405)

@csrf_exempt
@api_view(['GET', 'POST'])
@permission_classes([IsAuthenticated])
def profile(request):
    user = request.user
    logger.debug(f"User {user.username} accessed the profile view.")
    
    try:
        profile = Profile.objects.get(user=user)
    except Profile.DoesNotExist:
        logger.error(f"Profile for user {user.username} does not exist.")
        return JsonResponse({'error': 'Profile not found.'}, status=404)
    
    errors = []

    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            logger.debug(f"Received POST data: {data}")
        except json.JSONDecodeError:
            logger.error("Invalid JSON received.")
            return JsonResponse({'error': 'Invalid JSON.'}, status=400)

        # Update user and profile data based on received JSON
        user.first_name = data.get('first_name', user.first_name)
        user.last_name = data.get('last_name', user.last_name)
        profile.bio = data.get('bio', profile.bio)
        profile.location = data.get('location', profile.location)

        birth_date = data.get('birth_date')
        if birth_date:
            parsed_date = parse_date(birth_date)
            if parsed_date:
                profile.birth_date = parsed_date
                logger.debug(f"Updated birth date to {parsed_date}")
            else:
                errors.append("Invalid date format for birth date.")
                logger.warning("Invalid date format received for birth date.")
                profile.birth_date = None

        if not user.first_name:
            errors.append("First name is required.")
            logger.warning("First name is missing.")
        if not user.last_name:
            errors.append("Last name is required.")
            logger.warning("Last name is missing.")

        if not errors:
            user.save()
            profile.save()
            logger.info(f"Profile for user {user.username} updated successfully.")
            return JsonResponse({'message': 'Profile updated successfully.'})
        else:
            logger.error(f"Errors occurred: {errors}")
            return JsonResponse({'errors': errors}, status=400)

    response_data = {
        'user': {
            'first_name': user.first_name,
            'last_name': user.last_name,
        },
        'profile': {
            'bio': profile.bio,
            'location': profile.location,
            'birth_date': profile.birth_date.isoformat() if profile.birth_date else None
        }
    }

    logger.debug(f"Returning response data: {response_data}")
    return JsonResponse(response_data)


    response_data = {
        'user': {
            'firstName': user.firstName,
            'lastName': user.lastName,
        },
        'profile': {
            'bio': profile.bio,
            'location': profile.location,
            'birth_date': profile.birth_date.isoformat() if profile.birth_date else None
        }
    }

    logger.debug(f"Returning response data: {response_data}")
    return JsonResponse(response_data)


@csrf_exempt
@api_view(['POST'])
@permission_classes([IsAuthenticated])
def change_password(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            logger.debug(f"Received POST data: {data}")
            
            current_password = data.get('current_password')
            new_password = data.get('new_password')
            confirm_new_password = data.get('confirm_new_password')

            if not request.user.check_password(current_password):
                logger.warning(f"User {request.user.username} provided incorrect current password.")
                return JsonResponse({'error': 'Current password is incorrect.'}, status=400)

            if new_password != confirm_new_password:
                logger.warning(f"User {request.user.username} provided non-matching new passwords.")
                return JsonResponse({'error': 'New passwords do not match.'}, status=400)

            if new_password == current_password:
                logger.warning(f"User {request.user.username} attempted to use the same new password as current password.")
                return JsonResponse({'error': 'New password cannot be the same as the current password.'}, status=400)

            request.user.set_password(new_password)
            request.user.save()
            update_session_auth_hash(request, request.user)  # Important, to keep the user logged in
            logger.info(f"User {request.user.username} successfully changed their password.")
            return JsonResponse({'message': 'Password changed successfully.'})
        
        except json.JSONDecodeError:
            logger.error("Invalid JSON received.")
            return JsonResponse({'error': 'Invalid JSON.'}, status=400)

    logger.error("Invalid request method used.")
    return JsonResponse({'error': 'Invalid request method.'}, status=405)

@csrf_exempt
@api_view(['POST'])
@permission_classes([IsAuthenticated])
def summarize_document(request):
    try:
        # Extract and decrypt the incoming payload
        encrypted_content = request.POST.get('encrypted_content')
        if not encrypted_content:
            logger.warning("No encrypted content found in the request.")
            return JsonResponse({'error': 'No encrypted content found in the request.'}, status=400)

        logger.debug(f"Encrypted content received: {encrypted_content}")

        decrypted_content = decrypt_data(encrypted_content)
        logger.debug(f"Decrypted content: {decrypted_content}")

        data = json.loads(decrypted_content)

        document_context = data.get('documentContext')
        main_subject = data.get('mainSubject')
        summary_purpose = data.get('summary_purpose')
        length_detail = data.get('length_detail')
        important_elements = data.get('important_elements')
        audience = data.get('audience')
        tone = data.get('tone')
        format = data.get('format')
        additional_instructions = data.get('additional_instructions')
        document_content = data.get('document_content')  # Get document content

        logger.debug(f"Document parameters received: context={document_context}, subject={main_subject}, purpose={summary_purpose}")

        # Process the uploaded document
        document = request.FILES.get('document')
        if document:
            # Assuming the document is a text file
            document_content = document.read().decode('utf-8')
            logger.debug(f"Document content read from file.")
        
        if not document_content:
            logger.warning("No document content available.")
            return JsonResponse({'error': 'No document content available.'}, status=400)

        # Generate summary based on the document content and other parameters
        logger.info("Generating summary...")
        summary_content = generate_summary(
            document_context,
            main_subject,
            summary_purpose,
            length_detail,
            important_elements,
            audience,
            tone,
            format,
            additional_instructions,
            document_content  # Pass document content to the summary function
        )

        if summary_content:
            logger.info("Summary generated successfully.")
            encrypted_summary = encrypt_data({'generated_content': summary_content})
            return JsonResponse({'encrypted_content': encrypted_summary}, status=200)

        logger.error("Failed to generate summary.")
        return JsonResponse({'error': 'Failed to generate summary. Please try again.'}, status=500)

    except Exception as e:
        logger.error(f"An error occurred: {str(e)}")
        return JsonResponse({'error': str(e)}, status=500)


@csrf_exempt
@api_view(['POST'])
@permission_classes([IsAuthenticated])
def content_generator(request):
    try:
        # Load and decode the request body
        body = request.body.decode('utf-8')
        logger.debug(f"Request body received: {body}")

        # Extract and decrypt the incoming payload
        data = json.loads(body)
        encrypted_content = data.get('encrypted_content')
        if not encrypted_content:
            logger.warning("No encrypted content found in the request.")
            return JsonResponse({'error': 'No encrypted content found in the request.'}, status=400)

        logger.debug(f"Encrypted content received: {encrypted_content}")

        decrypted_content = decrypt_data(encrypted_content)
        logger.debug(f"Decrypted content: {decrypted_content}")

        data = json.loads(decrypted_content)

        # Extract fields from the decrypted JSON data
        company_info = data.get('company_info')
        content_purpose = data.get('content_purpose')
        desired_action = data.get('desired_action')
        topic_details = data.get('topic_details')
        keywords = data.get('keywords')
        audience_profile = data.get('audience_profile')
        format_structure = data.get('format_structure')
        num_words = data.get('num_words')
        seo_keywords = data.get('seo_keywords')
        references = data.get('references')

        logger.debug(f"Data extracted for content generation: company_info={company_info}, content_purpose={content_purpose}, desired_action={desired_action}")

        # Generate the content
        logger.info("Generating content...")
        content = generate_content(
            company_info,
            content_purpose,
            desired_action,
            topic_details,
            keywords,
            audience_profile,
            format_structure,
            num_words,
            seo_keywords,
            references
        )

        if content:
            logger.info("Content generated successfully.")
            encrypted_response_content = encrypt_data({'generated_content': content})
            return JsonResponse({'encrypted_content': encrypted_response_content}, status=200)

        logger.error("Failed to generate content.")
        return JsonResponse({'error': 'Failed to generate content. Please try again.'}, status=500)

    except json.JSONDecodeError:
        logger.error("Invalid JSON format received.")
        return JsonResponse({'error': 'Invalid JSON format. Please provide valid JSON data.'}, status=400)
    except ValueError as e:
        logger.error(f"ValueError occurred: {str(e)}")
        return JsonResponse({'error': str(e)}, status=400)
    except Exception as e:
        logger.error(f"An unexpected error occurred: {str(e)}")
        return JsonResponse({'error': str(e)}, status=500)

    logger.error("Method not allowed.")
    return JsonResponse({'error': 'Method not allowed.'}, status=405)

@csrf_exempt
@api_view(['POST'])
@permission_classes([IsAuthenticated])
def sales_script_generator(request):
    try:
        # Load and decode the request body
        body = request.body.decode('utf-8')
        logger.debug(f"Request body received: {body}")

        # Extract and decrypt the incoming payload
        data = json.loads(body)
        encrypted_content = data.get('encrypted_content')
        if not encrypted_content:
            logger.warning("No encrypted content found in the request.")
            return JsonResponse({'error': 'No encrypted content found in the request.'}, status=400)

        logger.debug(f"Encrypted content received: {encrypted_content}")

        decrypted_content = decrypt_data(encrypted_content)
        logger.debug(f"Decrypted content: {decrypted_content}")

        data = json.loads(decrypted_content)

        # Extract fields from the decrypted JSON data
        num_words = data.get('num_words')
        company_details = data.get('company_details')
        product_descriptions = data.get('product_descriptions')
        features_benefits = data.get('features_benefits')
        pricing_info = data.get('pricing_info')
        promotions = data.get('promotions')
        target_audience = data.get('target_audience')
        sales_objectives = data.get('sales_objectives')
        competitive_advantage = data.get('competitive_advantage')
        compliance = data.get('compliance')

        logger.debug(f"Data extracted for sales script generation: num_words={num_words}, company_details={company_details}")

        # Generate the sales script
        logger.info("Generating sales script...")
        sales_script = generate_sales_script(
            company_details,
            num_words,
            product_descriptions,
            features_benefits,
            pricing_info,
            promotions,
            target_audience,
            sales_objectives,
            competitive_advantage,
            compliance,
        )

        if sales_script:
            logger.info("Sales script generated successfully.")
            encrypted_response_content = encrypt_data({'generated_content': sales_script})
            return JsonResponse({'encrypted_content': encrypted_response_content}, status=200)

        logger.error("Failed to generate sales script.")
        return JsonResponse({'error': 'Failed to generate sales script. Please try again.'}, status=500)

    except json.JSONDecodeError:
        logger.error("Invalid JSON format received.")
        return JsonResponse({'error': 'Invalid JSON format. Please provide valid JSON data.'}, status=400)
    except ValueError as e:
        logger.error(f"ValueError occurred: {str(e)}")
        return JsonResponse({'error': str(e)}, status=400)
    except Exception as e:
        logger.error(f"An unexpected error occurred: {str(e)}")
        return JsonResponse({'error': str(e)}, status=500)

    logger.error("Method not allowed.")
    return JsonResponse({'error': 'Method not allowed.'}, status=405)

@csrf_exempt

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def logout_view(request):
    try:
        # Perform logout operation
        logout(request)
        logger.info(f"User {request.user.username} logged out successfully.")
        return JsonResponse({'success': 'Logged out successfully'}, status=200)
    except Exception as e:
        logger.error(f"Error during logout: {str(e)}")
        return JsonResponse({'error': 'An error occurred during logout.'}, status=500)

# @csrf_exempt
# @api_view(['POST'])
# @permission_classes([IsAuthenticated])
# def create_presentation(request):
#     try:
#         # Load and parse the JSON data
#         data = json.loads(request.POST.get('data'))
#         logger.debug(f"Data received: {data}")

#         title = data.get('title')
#         num_slides = data.get('num_slides')
#         special_instructions = data.get('special_instructions')
        
#         # Process background image file
#         bg_image_file = request.FILES.get('bg_image', None)
#         if bg_image_file:
#             bg_image_path = default_storage.save(bg_image_file.name, bg_image_file)
#             bg_image = default_storage.path(bg_image_path)
#             logger.debug(f"Background image saved at: {bg_image}")
#         else:
#             bg_image = None

#         # Process document file
#         document_file = request.FILES.get('document', None)
#         if document_file:
#             document_path = default_storage.save(document_file.name, document_file)
#             document_content = extract_document_content(default_storage.path(document_path))
#             logger.debug("Document content extracted.")
#         else:
#             document_content = None

#         # Generate presentation
#         logger.info("Generating presentation...")
#         prs = Presentation()
#         slide_titles = generate_slide_titles(title, num_slides, special_instructions)
#         slide_titles = slide_titles.replace('[', '').replace(']', '').replace('"', '').split(',')

#         max_points_per_slide = 4

#         for st in slide_titles:
#             slide_content = generate_slide_content(st, title, special_instructions, document_content).replace("*", '').split('\n')
#             current_content = []
#             slide_count = 1

#             for point in slide_content:
#                 current_content.append(point.strip())
#                 if len(current_content) >= max_points_per_slide:
#                     add_slide(prs, st if slide_count == 1 else f"{st} (contd.)", current_content, bg_image)
#                     current_content = []
#                     slide_count += 1

#             if current_content:
#                 add_slide(prs, st if slide_count == 1 else f"{st} (contd.)", current_content, bg_image)

#         # Save presentation to buffer and prepare response
#         pptx_buffer = BytesIO()
#         prs.save(pptx_buffer)
#         pptx_buffer.seek(0)

#         response = HttpResponse(pptx_buffer, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
#         response['Content-Disposition'] = 'attachment; filename="SmartOffice_Assistant_Presentation.pptx"'

#         logger.info("Presentation generated and response prepared.")
#         return response

#     except json.JSONDecodeError:
#         logger.error("Invalid JSON format received.")
#         return JsonResponse({'error': 'Invalid JSON format. Please provide valid JSON data.'}, status=400)
#     except ValueError as e:
#         logger.error(f"ValueError occurred: {str(e)}")
#         return JsonResponse({'error': str(e)}, status=400)
#     except Exception as e:
#         logger.error(f"An unexpected error occurred: {str(e)}")
#         return JsonResponse({'error': str(e)}, status=500)

#     logger.error("Method not allowed.")
#     return JsonResponse({'error': 'Method not allowed.'}, status=405)


@csrf_exempt
@api_view(['POST'])
@permission_classes([IsAuthenticated])
def create_presentation(request):
    try:
        # Load and parse the JSON data
        data = json.loads(request.POST.get('data'))
        logger.debug(f"Data received: {data}")

        title = data.get('title')
        num_slides = int(data.get('num_slides'))  # Ensure num_slides is an integer
        special_instructions = data.get('special_instructions')
        
        # Process background image file
        bg_image_file = request.FILES.get('bg_image', None)
        if bg_image_file:
            bg_image_path = default_storage.save(bg_image_file.name, bg_image_file)
            bg_image = default_storage.path(bg_image_path)
            logger.debug(f"Background image saved at: {bg_image}")
        else:
            bg_image = None

        # Process document file
        document_file = request.FILES.get('document', None)
        if document_file:
            document_path = default_storage.save(document_file.name, document_file)
            document_content = extract_document_content(default_storage.path(document_path))
            logger.debug("Document content extracted.")
        else:
            document_content = None

        # Generate presentation
        logger.info("Generating presentation...")
        prs = Presentation()
        slide_titles = generate_slide_titles(title, num_slides, special_instructions)
        slide_titles = slide_titles.replace('[', '').replace(']', '').replace('"', '').split(',')

        max_points_per_slide = 4
        total_slides_generated = 0

        for st in slide_titles:
            if total_slides_generated >= num_slides:
                break

            slide_content = generate_slide_content(st, title, special_instructions, document_content).replace("*", '').split('\n')
            current_content = []
            slide_count = 1

            for point in slide_content:
                current_content.append(point.strip())
                if len(current_content) >= max_points_per_slide:
                    add_slide(prs, st if slide_count == 1 else f"{st} (contd.)", current_content, bg_image)
                    current_content = []
                    slide_count += 1
                    total_slides_generated += 1

                    if total_slides_generated >= num_slides:
                        break

            if current_content and total_slides_generated < num_slides:
                add_slide(prs, st if slide_count == 1 else f"{st} (contd.)", current_content, bg_image)
                total_slides_generated += 1

        # Save presentation to buffer and prepare response
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        response = HttpResponse(pptx_buffer, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="SmartOffice_Assistant_Presentation.pptx"'

        logger.info("Presentation generated and response prepared.")
        return response

    except json.JSONDecodeError:
        logger.error("Invalid JSON format received.")
        return JsonResponse({'error': 'Invalid JSON format. Please provide valid JSON data.'}, status=400)
    except ValueError as e:
        logger.error(f"ValueError occurred: {str(e)}")
        return JsonResponse({'error': str(e)}, status=400)
    except Exception as e:
        logger.error(f"An unexpected error occurred: {str(e)}")
        return JsonResponse({'error': str(e)}, status=500)

    logger.error("Method not allowed.")
    return JsonResponse({'error': 'Method not allowed.'}, status=405)